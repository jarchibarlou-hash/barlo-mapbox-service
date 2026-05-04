#!/usr/bin/env python3
"""
BARLO -- Server-side PPTX generation from diagnostic data.
v4.0 -- Premium density: font sizes calibrated, structural text templates.
"""
import json, sys, os, re, copy, tempfile, urllib.request, shutil
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from generate_charts import generate_all_charts

# -------------------------------------------------------------
# PLACEHOLDER MAPPINGS
# -------------------------------------------------------------

IMAGE_PLACEHOLDERS = {
    '{{slide_4_image}}',
    '{{slide_4_axo_image}}',
    '{{scenario_A_massing}}',
    '{{scenario_B_massing}}',
    '{{scenario_C_massing}}',
}

CHART_PLACEHOLDERS = {
    '{{scenario_A_risk_chart}}': ['scenario_A_risk_radar', 'scenario_A_risk_gauge', 'scenario_A_risk_bars'],
    '{{scenario_B_risk_chart}}': ['scenario_B_risk_radar', 'scenario_B_risk_gauge', 'scenario_B_risk_bars'],
    '{{scenario_C_risk_chart}}': ['scenario_C_risk_radar', 'scenario_C_risk_gauge', 'scenario_C_risk_bars'],
    '{{tableau comparative_charts}}': ['tableau_comparative_charts'],
    '{{arbitrage_graph_}}': ['arbitrage_graph_'],
}

SLIDE_SPECIFIC_TEXT = {
    (17, '{{invisible_technical_text}}'): 'invisible_technical_text_s17',
    (17, '{{invisible_financial_text}}'): 'invisible_financial_text_s17',
    (17, '{{invisible_strategic_text}}'): 'invisible_strategic_text_s17',
    (18, '{{invisible_technical_text}}'): 'invisible_technical_text_s18',
    (18, '{{invisible_financial_text}}'): 'invisible_financial_text_s18',
    (18, '{{invisible_strategic_text}}'): 'invisible_strategic_text_s18',
}

# Risk chart slides that need fallback shape detection
RISK_CHART_SLIDES = {8, 11, 14}

# -------------------------------------------------------------
# UTILITY FUNCTIONS
# -------------------------------------------------------------

def download_image(url, dest_dir):
    try:
        filename = os.path.join(dest_dir, os.path.basename(url).split('?')[0] or 'image.png')
        req = urllib.request.Request(url, headers={'User-Agent': 'BARLO-PPTX/1.0'})
        with urllib.request.urlopen(req, timeout=30) as resp:
            with open(filename, 'wb') as f:
                shutil.copyfileobj(resp, f)
        return filename
    except Exception as e:
        print(f"WARNING: Failed to download {url}: {e}", file=sys.stderr)
        return None

def find_placeholder_in_shape(shape, placeholder_text):
    if not shape.has_text_frame:
        return False
    full_text = ''.join(run.text for para in shape.text_frame.paragraphs for run in para.runs)
    if not full_text:
        full_text = shape.text_frame.text
    return placeholder_text in full_text

def get_shape_placeholder(shape):
    if not shape.has_text_frame:
        return None
    full_text = shape.text_frame.text
    match = re.search(r'\{\{[^}]+\}\}', full_text)
    if match:
        return match.group(0)
    match = re.search(r'\{\{([a-zA-Z_][a-zA-Z0-9_ ]*)\}', full_text)
    if match:
        return '{{' + match.group(1) + '}}'
    return None

def clear_shape_text(shape):
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = ''

# -------------------------------------------------------------
# FONT SIZE STRATEGY
# -------------------------------------------------------------
# Template uses 15pt body / 16pt client_name.
# Reference PDF shows ~13-14pt for body text (auto-shrunk from 15pt).
# We keep template fonts for most slides and only override for dense slides.
# fontScale=80000 (80%) prevents over-shrinking.

def get_font_size_for_slide(slide_num):
    """
    Return target font size matching premium document density.
    v74.13 -- Calibration : textes refondus plus courts (lot P4) → on peut
    remonter les tailles. Cible architecte conseil moderne, lisible.
    Slides 8/11/14/16 spécifiquement remontées (feedback user "trop petits").
    """
    if slide_num == 3:
        return Pt(12)    # Intro -- moderate density, slightly bumped
    elif slide_num in [6, 9, 12]:
        return Pt(11)    # Scenario summaries -- bumped from 9 (textes raccourcis)
    elif slide_num == 5:
        return Pt(11)    # Context -- bumped from 9
    elif slide_num in [8, 11, 14]:
        return Pt(12)    # Risk slides -- BUMPED from 10 (feedback "trop petits")
    elif slide_num in [4, 7, 10, 13, 19]:
        return Pt(11)    # Site, financial, next steps -- bumped from 10
    elif slide_num == 16:
        return Pt(12)    # Strategic arbitrage -- BUMPED from 10 (feedback "trop petit")
    elif slide_num in [17, 18]:
        return Pt(10)    # 3-column layout -- bumped from 9 (still dense but legible)
    elif slide_num == 20:
        return Pt(12)    # Conclusion -- bumped from 11
    else:
        return None       # Keep template font (15-16pt) for slides 1, 2, 15

def set_font_size_for_shape(shape, font_size_pt):
    """
    Set all text runs in a shape to a specific font size.
    Only called when font_size_pt is not None.
    """
    if not shape.has_text_frame or font_size_pt is None:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = font_size_pt

def enable_auto_shrink(shape, fontScale=80000):
    """
    Enable auto-shrink for text overflow.
    fontScale=80000 = 80% minimum (15pt base → 12pt min, 12pt base → 9.6pt min)
    """
    if not shape.has_text_frame:
        return
    from lxml import etree
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    txBody = shape.text_frame._txBody
    bodyPr = txBody.find(f'{{{ns}}}bodyPr')
    if bodyPr is None:
        return
    for child_tag in ['noAutofit', 'normAutofit', 'spAutoFit']:
        existing = bodyPr.find(f'{{{ns}}}{child_tag}')
        if existing is not None:
            bodyPr.remove(existing)
    normAutofit = etree.SubElement(bodyPr, f'{{{ns}}}normAutofit')
    normAutofit.set('fontScale', str(fontScale))

# -------------------------------------------------------------
# TEXT REPLACEMENT
# -------------------------------------------------------------

def replace_text_in_shape(shape, placeholder, new_text):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    placeholder_found = False
    placeholder_core = placeholder.strip('{}')
    for para in tf.paragraphs:
        para_text = ''.join(run.text for run in para.runs)
        if not para_text:
            para_text = para.text
        if placeholder in para_text or ('{{' + placeholder_core) in para_text:
            placeholder_found = True
            break
    if not placeholder_found:
        return
    # Get reference run for formatting
    ref_run = None
    for para in tf.paragraphs:
        for run in para.runs:
            ref_run = run
            break
        if ref_run:
            break
    new_paragraphs = new_text.split('\n') if new_text else ['']
    xml_element = tf._txBody
    p_elements = xml_element.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}p')
    for p_elem in p_elements[1:]:
        xml_element.remove(p_elem)
    first_p = p_elements[0]
    r_elements = first_p.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}r')
    for r_elem in r_elements:
        first_p.remove(r_elem)
    _add_run_to_paragraph(first_p, new_paragraphs[0], ref_run)
    for para_text in new_paragraphs[1:]:
        new_p = copy.deepcopy(first_p)
        r_elements = new_p.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}r')
        for r_elem in r_elements:
            new_p.remove(r_elem)
        _add_run_to_paragraph(new_p, para_text, ref_run)
        xml_element.append(new_p)

def _parse_markdown_bold(text):
    """v74.12 — Parse markdown **bold** into [(text, is_bold), ...] segments.
    Robust against unbalanced markers : ** sans pair → traite comme texte normal."""
    if not text or '**' not in text:
        return [(text or '', False)]
    segments = []
    parts = text.split('**')
    # Si nombre impair de **, le dernier segment reste en non-bold (markers déséquilibrés)
    is_bold = False
    for part in parts:
        if part:
            segments.append((part, is_bold))
        is_bold = not is_bold
    # Si on a fini en is_bold=True (impair), le dernier ** orphelin a inversé pour rien.
    # Pas de souci, le segment a déjà été ajouté avec son flag correct.
    return segments

def _add_run_to_paragraph(p_element, text, ref_run=None):
    """v74.12 — Crée 1+ runs en parsant markdown **bold**.
    Si pas de **, comportement identique à avant (1 seul run).
    Sinon, alterne runs normaux et runs bold."""
    from lxml import etree
    A_NS = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    segments = _parse_markdown_bold(text)
    # Reference rPr depuis ref_run (formatting hérité) — sans le b
    base_rPr = None
    if ref_run is not None and ref_run._r is not None:
        rPr_orig = ref_run._r.find(f'{A_NS}rPr')
        if rPr_orig is not None:
            base_rPr = copy.deepcopy(rPr_orig)
            base_rPr.attrib.pop('b', None)
    for seg_text, seg_bold in segments:
        if not seg_text:
            continue
        r = etree.SubElement(p_element, f'{A_NS}r')
        if base_rPr is not None:
            seg_rPr = copy.deepcopy(base_rPr)
            if seg_bold:
                seg_rPr.set('b', '1')
            r.insert(0, seg_rPr)
        elif seg_bold:
            seg_rPr = etree.SubElement(r, f'{A_NS}rPr')
            seg_rPr.set('b', '1')
        t = etree.SubElement(r, f'{A_NS}t')
        t.text = seg_text
        if seg_text and (seg_text[0] == ' ' or seg_text[-1] == ' '):
            t.set('{http://www.w3.org/XML/1998/namespace}space', 'preserve')

# -------------------------------------------------------------
# IMAGE REPLACEMENT
# -------------------------------------------------------------

def replace_shape_with_image(slide, shape, image_path, override_bounds=None, maintain_aspect_ratio=False):
    if not os.path.exists(image_path):
        print(f"WARNING: Image not found: {image_path}", file=sys.stderr)
        return
    if override_bounds:
        left, top, width, height = override_bounds
    else:
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height
    sp_element = shape._element
    sp_element.getparent().remove(sp_element)
    if maintain_aspect_ratio:
        try:
            from PIL import Image
            img = Image.open(image_path)
            img_width, img_height = img.size
            aspect_ratio = img_width / img_height if img_height > 0 else 1.0
            target_width = width
            target_height = int(target_width / aspect_ratio)
            if target_height > height:
                target_height = height
                target_width = int(target_height * aspect_ratio)
            width = target_width
            height = target_height
        except Exception as e:
            print(f"WARNING: Could not load image for aspect ratio: {e}", file=sys.stderr)
    slide.shapes.add_picture(image_path, left, top, width, height)

def replace_shape_with_multiple_images(slide, shape, image_paths, maintain_aspect_ratio=True):
    """v74.13 — par défaut maintain_aspect_ratio=True pour eviter la distorsion
    des charts (radar, barres, jauge) inseres cote-a-cote. Calcule pour chaque
    image la taille qui rentre dans la cellule (img_width × height) en preservant
    le ratio source, puis centre verticalement dans la cellule disponible."""
    valid_paths = [p for p in image_paths if p and os.path.exists(p)]
    if not valid_paths:
        print("WARNING: No valid images for multi-image replacement", file=sys.stderr)
        return
    left = shape.left
    top = shape.top
    total_width = shape.width
    height = shape.height
    n = len(valid_paths)
    gap = Emu(36000)
    img_width = (total_width - gap * (n - 1)) // n
    sp_element = shape._element
    sp_element.getparent().remove(sp_element)
    for i, img_path in enumerate(valid_paths):
        img_left = left + i * (img_width + gap)
        cell_w = img_width
        cell_h = height
        cell_top = top
        if maintain_aspect_ratio:
            try:
                from PIL import Image
                pil = Image.open(img_path)
                src_w, src_h = pil.size
                if src_h > 0 and src_w > 0:
                    ratio = src_w / src_h
                    target_w = cell_w
                    target_h = int(target_w / ratio)
                    if target_h > cell_h:
                        target_h = cell_h
                        target_w = int(target_h * ratio)
                    # Centrer verticalement dans la cellule disponible
                    cell_top = top + (cell_h - target_h) // 2
                    img_left = img_left + (cell_w - target_w) // 2
                    cell_w = target_w
                    cell_h = target_h
            except Exception as e:
                print(f"WARNING: Could not preserve aspect ratio for {img_path}: {e}", file=sys.stderr)
        slide.shapes.add_picture(img_path, img_left, cell_top, cell_w, cell_h)

# -------------------------------------------------------------
# SLIDE 15 -- COMPARATIVE TABLE (FULL SLIDE)
# -------------------------------------------------------------

SLIDE_15_GRID_SHAPES = {
    'Google Shape;141;p27', 'Google Shape;145;p27',
    'Google Shape;146;p27', 'Google Shape;147;p27',
    'Google Shape;150;p27', 'Google Shape;151;p27',
    'Google Shape;152;p27', 'Google Shape;153;p27',
}

def _handle_slide_15(slide, chart_paths):
    chart_path = chart_paths.get('tableau_comparative_charts')
    if not chart_path or not os.path.exists(chart_path):
        print("WARNING: No comparatif chart for slide 15", file=sys.stderr)
        return
    shapes_to_remove = []
    for shape in slide.shapes:
        name = shape.name if hasattr(shape, 'name') else ''
        placeholder = get_shape_placeholder(shape)
        if placeholder == '{{tableau comparative_charts}}' or name in SLIDE_15_GRID_SHAPES:
            shapes_to_remove.append(shape)
    for shape in shapes_to_remove:
        sp_element = shape._element
        sp_element.getparent().remove(sp_element)
    # Full width with small margins
    left = Emu(200000)
    top = Emu(391320)
    width = Emu(8700000)
    height = Emu(4400000)
    slide.shapes.add_picture(chart_path, left, top, width, height)

# -------------------------------------------------------------
# RISK CHART FALLBACK (SLIDES 8, 11, 14)
# -------------------------------------------------------------

def find_large_shape_for_charts(slide):
    """Find large empty shape on risk chart slides."""
    min_width = Emu(914400 * 8)
    min_height = Emu(914400 * 2)    # Relaxed from 2.5 to 2 inches
    max_top = Emu(914400 * 2)       # Relaxed from 1.5 to 2 inches

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text_content = shape.text_frame.text.strip()
        if len(text_content) > 50:
            continue
        if shape.width >= min_width and shape.height >= min_height and shape.top <= max_top:
            return shape
    return None

def _insert_risk_charts_fallback(slide, slide_num, chart_paths):
    """Fallback: find large shape and insert 3 risk charts side-by-side."""
    scenario_map = {8: 'scenario_A', 11: 'scenario_B', 14: 'scenario_C'}
    scenario = scenario_map.get(slide_num)
    if not scenario:
        return False
    chart_keys = [f'{scenario}_risk_radar', f'{scenario}_risk_gauge', f'{scenario}_risk_bars']
    chart_image_paths = [chart_paths.get(k) for k in chart_keys]
    chart_image_paths = [p for p in chart_image_paths if p and os.path.exists(p)]
    if not chart_image_paths:
        print(f"WARNING: No risk charts for slide {slide_num}", file=sys.stderr)
        return False
    target_shape = find_large_shape_for_charts(slide)
    if not target_shape:
        print(f"WARNING: No suitable shape found for risk charts on slide {slide_num}", file=sys.stderr)
        return False
    print(f"Found target shape for risk charts on slide {slide_num}: {target_shape.name}", file=sys.stderr)
    if len(chart_image_paths) == 1:
        # v74.13 : preserve aspect ratio sur charts (homothecie OK)
        replace_shape_with_image(slide, target_shape, chart_image_paths[0], maintain_aspect_ratio=True)
    else:
        replace_shape_with_multiple_images(slide, target_shape, chart_image_paths, maintain_aspect_ratio=True)
    return True

# -------------------------------------------------------------
# PHASAGE TEXT CLEANUP (SLIDE 18)
# -------------------------------------------------------------

def _clean_phasage_text(text):
    """
    Remove raw phasage data from text.
    Phasage info is displayed as a timeline chart on slide 19 instead.
    """
    if not text:
        return ''
    text_upper = text.upper()
    if 'PHASAGE' in text_upper or 'MONOPHASEE' in text_upper or 'TRIPHASEE' in text_upper or 'BIPHASEE' in text_upper or 'MONOPHASÉE' in text_upper:
        return ''
    return text

# -------------------------------------------------------------
# APPLY TEXT + FONT TO A SHAPE (DRY helper)
# -------------------------------------------------------------

def _apply_text_to_shape(shape, placeholder, text, slide_num):
    """v74.14 — calibration auto-shrink :
    - Slides 8/11/14/16 (textes refondus courts) : 95% min, presque pas de shrink → on garde les 12pt
    - Slides 5/6/9/12 : 85% min (un peu plus serre)
    - Slides 17/18 (3-col tres dense) : 80% min
    - Autres : 90% min standard
    """
    replace_text_in_shape(shape, placeholder, text)
    font_size = get_font_size_for_slide(slide_num)
    if font_size is not None:
        set_font_size_for_shape(shape, font_size)
    if slide_num in [8, 11, 14, 16]:
        enable_auto_shrink(shape, fontScale=95000)  # 95% min — quasi pas de shrink, garantit la 12pt
    elif slide_num in [17, 18]:
        enable_auto_shrink(shape, fontScale=80000)  # 80% — 3-col dense
    elif slide_num in [5, 6, 9, 12]:
        enable_auto_shrink(shape, fontScale=85000)  # 85% — scenario sommaires
    else:
        enable_auto_shrink(shape, fontScale=90000)  # 90% standard (vs 80% avant)

# -------------------------------------------------------------
# MAIN ASSEMBLY
# -------------------------------------------------------------

def assemble_pptx(data, template_path, output_path):
    chart_dir = tempfile.mkdtemp(prefix='barlo_charts_')
    print(f"Generating charts in {chart_dir}...", file=sys.stderr)
    chart_paths = generate_all_charts(data, chart_dir)
    print(f"Charts generated: {list(chart_paths.keys())}", file=sys.stderr)

    img_dir = tempfile.mkdtemp(prefix='barlo_images_')
    images = data.get('images', {})
    downloaded_images = {}
    for key, url in images.items():
        if url:
            print(f"Downloading image: {key}...", file=sys.stderr)
            local_path = download_image(url, img_dir)
            if local_path:
                downloaded_images[key] = local_path

    prs = Presentation(template_path)
    texts = data.get('texts', {})
    client_name = data.get('client_name', '')

    # -- Prepare composite text keys --
    if 'slide_3_text' not in texts and 'slide_3_intro_text' in texts:
        parts = [texts.get('slide_3_intro_text', ''), texts.get('slide_3_programme_text', '')]
        texts['slide_3_text'] = '\n\n'.join(p for p in parts if p)

    for base_key in ['invisible_technical_text', 'invisible_financial_text', 'invisible_strategic_text']:
        s17_key = f'{base_key}_s17'
        s18_key = f'{base_key}_s18'
        if s17_key not in texts and base_key in texts:
            texts[s17_key] = texts[base_key]
        if s18_key not in texts:
            success_key = base_key.replace('invisible_', 'success_')
            if success_key in texts:
                texts[s18_key] = texts[success_key]
            elif base_key in texts:
                texts[s18_key] = texts[base_key]

    print(f"Text keys available: {sorted(texts.keys())}", file=sys.stderr)
    print(f"Chart paths available: {sorted(chart_paths.keys())}", file=sys.stderr)

    # -- Process each slide --
    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1

        # Slide 15 -- full-page comparative table
        if slide_num == 15:
            _handle_slide_15(slide, chart_paths)
            continue

        # Risk chart slides (8, 11, 14) -- special chart insertion + text
        if slide_num in RISK_CHART_SLIDES:
            shapes_to_process = list(slide.shapes)
            chart_inserted = False
            for shape in shapes_to_process:
                placeholder = get_shape_placeholder(shape)
                if placeholder and placeholder in CHART_PLACEHOLDERS:
                    chart_keys = CHART_PLACEHOLDERS[placeholder]
                    chart_image_paths = [chart_paths.get(k) for k in chart_keys]
                    chart_image_paths = [p for p in chart_image_paths if p]
                    if len(chart_image_paths) == 1:
                        # v74.13 : preserve aspect ratio sur charts
                        replace_shape_with_image(slide, shape, chart_image_paths[0], maintain_aspect_ratio=True)
                        chart_inserted = True
                    elif len(chart_image_paths) > 1:
                        replace_shape_with_multiple_images(slide, shape, chart_image_paths, maintain_aspect_ratio=True)
                        chart_inserted = True
                    break
            if not chart_inserted:
                print(f"No risk chart placeholder found on slide {slide_num}, using fallback...", file=sys.stderr)
                _insert_risk_charts_fallback(slide, slide_num, chart_paths)

            # Process text shapes on risk slides
            shapes_to_process = list(slide.shapes)
            for shape in shapes_to_process:
                placeholder = get_shape_placeholder(shape)
                if not placeholder or placeholder in CHART_PLACEHOLDERS:
                    continue
                if placeholder == '{{client_name}}':
                    _apply_text_to_shape(shape, placeholder, client_name, slide_num)
                    continue
                placeholder_key = placeholder.strip('{}').strip()
                slide_specific_key = (slide_num, placeholder)
                if slide_specific_key in SLIDE_SPECIFIC_TEXT:
                    text_key = SLIDE_SPECIFIC_TEXT[slide_specific_key]
                    text = texts.get(text_key, '')
                    if text:
                        _apply_text_to_shape(shape, placeholder, text, slide_num)
                    else:
                        clear_shape_text(shape)
                    continue
                text = texts.get(placeholder_key, '')
                if text:
                    _apply_text_to_shape(shape, placeholder, text, slide_num)
                else:
                    clean_key = placeholder_key.replace(' ', '_')
                    text = texts.get(clean_key, '')
                    if text:
                        _apply_text_to_shape(shape, placeholder, text, slide_num)

        else:
            # Standard shape processing for all non-risk, non-slide-15 slides
            shapes_to_process = list(slide.shapes)
            for shape in shapes_to_process:
                placeholder = get_shape_placeholder(shape)
                if not placeholder:
                    continue
                placeholder_key = placeholder.strip('{}').strip()

                # Client name
                if placeholder == '{{client_name}}':
                    replace_text_in_shape(shape, placeholder, client_name)
                    # Slide 1: keep template 16pt, no font override
                    # Other slides: apply slide-specific font
                    if slide_num != 1:
                        font_size = get_font_size_for_slide(slide_num)
                        if font_size is not None:
                            set_font_size_for_shape(shape, font_size)
                    enable_auto_shrink(shape, fontScale=80000)
                    continue

                # Images
                if placeholder in IMAGE_PLACEHOLDERS:
                    img_key = placeholder_key
                    img_path = downloaded_images.get(img_key)
                    if img_path:
                        # v73.2.1 : maintain_aspect_ratio=True pour TOUTES les images
                        # (axo brute slide 4, axo enhanced slide 5, massings A/B/C)
                        # Evite la deformation horizontale (homothetie respectee)
                        replace_shape_with_image(slide, shape, img_path, maintain_aspect_ratio=True)
                    else:
                        clear_shape_text(shape)
                    continue

                # Charts
                if placeholder in CHART_PLACEHOLDERS:
                    chart_keys = CHART_PLACEHOLDERS[placeholder]
                    chart_image_paths = [chart_paths.get(k) for k in chart_keys]
                    chart_image_paths = [p for p in chart_image_paths if p]
                    if len(chart_image_paths) == 1:
                        # v74.13 : preserve aspect ratio sur charts
                        replace_shape_with_image(slide, shape, chart_image_paths[0], maintain_aspect_ratio=True)
                    elif len(chart_image_paths) > 1:
                        replace_shape_with_multiple_images(slide, shape, chart_image_paths, maintain_aspect_ratio=True)
                    else:
                        clear_shape_text(shape)
                    continue

                # Slide-specific text (slides 17, 18)
                slide_specific_key = (slide_num, placeholder)
                if slide_specific_key in SLIDE_SPECIFIC_TEXT:
                    text_key = SLIDE_SPECIFIC_TEXT[slide_specific_key]
                    text = texts.get(text_key, '')
                    # v72.94: REMOVED phasage cleanup — _s18 texts are clean
                    # buildTemplateTexts already produces proper prose, not raw data
                    if text:
                        _apply_text_to_shape(shape, placeholder, text, slide_num)
                    else:
                        clear_shape_text(shape)
                    continue

                # Generic text
                text = texts.get(placeholder_key, '')
                if text:
                    _apply_text_to_shape(shape, placeholder, text, slide_num)
                else:
                    clean_key = placeholder_key.replace(' ', '_')
                    text = texts.get(clean_key, '')
                    if text:
                        _apply_text_to_shape(shape, placeholder, text, slide_num)
                    else:
                        print(f"WARNING: No text for placeholder {placeholder} on slide {slide_num}", file=sys.stderr)

    # ----------------------------------------------------------
    # INSERT EXTRA CHARTS (positioned within slide bounds)
    # Slide dimensions: 10.00" x 5.62" (9144000 x 5143500 EMU)
    # ----------------------------------------------------------

    slides_list = list(prs.slides)

    # Slide 17 -- Cost breakdown pie chart
    # v72.92: Repositioned to bottom-right to avoid overlapping text columns
    cost_chart = chart_paths.get('cost_breakdown')
    if cost_chart and os.path.exists(cost_chart) and len(slides_list) >= 17:
        slide17 = slides_list[16]
        # Right side, below text columns: left=6.3", top=3.3", width=3.2", height=2.1"
        slide17.shapes.add_picture(cost_chart,
            Emu(5760720), Emu(3017520), Emu(2926080), Emu(1920240))
        print("Inserted cost breakdown chart on slide 17 (bottom-right)", file=sys.stderr)

    # ─── v74.17 SLIDES 7/10/13 — DISCIPLINE LAYOUT ────────────────────────────
    # v25 a montre que le shape TEXTE prend toute la slide (~5.5" haut),
    # donc le texte affiche son contenu jusqu'a ~4.2" et SE CHEVAUCHE
    # avec mes charts a 3.0" et 4.1". Fix : RESIZE le shape texte
    # programmatiquement a 2.5" haut → 0.5"-3.0" → libere 3.0"-5.5" net.
    #
    # Plus : calc figsize=(11,3) ratio 3.67 inserer dans 9x1 ratio 9 → 2.4x
    # de stretch horizontal. v74.17 : ajuste figsize ET insertion pour
    # matcher 1:1 (zero distorsion).
    #
    # Layout cible final (zero overlap, zero distorsion) :
    #   - Title (template) : 0-0.5"
    #   - Texte (resize) : 0.5"-3.0" (h=2.5")
    #   - Calc visuel : 3.1"-4.4" (h=1.3", figsize 11×1.6 ratio 6.9 → insert 9.0×1.3 ratio 6.9 ✓)
    #   - Gauge budget : 4.5"-5.5" (h=1.0", figsize 11×1.4 ratio 7.9 → insert 9.0×1.0 ratio 9.0 close)
    FINANCIAL_SLIDE_MAP = {7: 'A', 10: 'B', 13: 'C'}
    for slide_num_fin, label_fin in FINANCIAL_SLIDE_MAP.items():
        if len(slides_list) < slide_num_fin:
            continue
        slide_fin = slides_list[slide_num_fin - 1]
        # ─── ETAPE 1 : Resize le shape texte principal (le plus grand non-titre) ──
        text_shapes = []
        for shp in slide_fin.shapes:
            try:
                if not shp.has_text_frame:
                    continue
                if shp.top is None or shp.top < Emu(700000):  # skip title (top<0.77")
                    continue
                text_shapes.append((shp.height, shp))
            except Exception:
                continue
        if text_shapes:
            text_shapes.sort(key=lambda t: t[0], reverse=True)
            main_text = text_shapes[0][1]
            try:
                main_text.height = Emu(2286000)  # 2.5"
                main_text.top = Emu(457200)      # 0.5"
                print(f"v74.17 slide {slide_num_fin}: text shape resized to top=0.5\" h=2.5\"", file=sys.stderr)
            except Exception as e:
                print(f"v74.17 slide {slide_num_fin}: failed to resize text shape: {e}", file=sys.stderr)
        # ─── ETAPE 2 : Inserer les 2 charts avec espace entre eux ──
        gauge = chart_paths.get(f'scenario_{label_fin}_budget_gauge')
        calc = chart_paths.get(f'scenario_{label_fin}_cost_calc')
        # Calc visuel : top=3.0", left=0.5", w=9.0", h=1.1"
        if calc and os.path.exists(calc):
            slide_fin.shapes.add_picture(calc,
                Emu(457200), Emu(2743200),  # left=0.5", top=3.0"
                Emu(8229600), Emu(1005840)) # 9.0" × 1.1"
        # Gauge budget : top=4.4", left=0.5", w=9.0", h=1.0" (gap 0.3" entre les deux)
        if gauge and os.path.exists(gauge):
            slide_fin.shapes.add_picture(gauge,
                Emu(457200), Emu(4023360),  # left=0.5", top=4.4"
                Emu(8229600), Emu(914400))  # 9.0" × 1.0"
        print(f"v74.18: Inserted financial charts on slide {slide_num_fin} ({label_fin})", file=sys.stderr)

    # Slide 17 -- Budget comparison table (bottom-left, next to pie chart)
    # v72.92: Repositioned to bottom-left to coexist with pie chart on right
    # Table: Scenario | SDP | Cout/m2 marche | Cout/m2 ajuste | Cout total | Label
    if len(slides_list) >= 17:
        slide17 = slides_list[16]
        try:
            # Extract data from flat_data
            table_rows = []
            for sc_key in ['A', 'B', 'C']:
                sdp_val = flat_data.get(f'{sc_key}_sdp', '0')
                cost_m2_marche = flat_data.get(f'{sc_key}_cost_m2_marche', '0k')
                cost_m2_ajuste = flat_data.get(f'{sc_key}_cost_m2_ajuste', '0k')
                cost_total = flat_data.get(f'{sc_key}_cost_total', '0M FCFA')
                budget_fit = flat_data.get(f'{sc_key}_budget_fit', '')
                # Translate budget_fit labels
                fit_label = {
                    'DANS_BUDGET': 'DANS BUDGET',
                    'BUDGET_TENDU': 'BUDGET TENDU',
                    'HORS_BUDGET': 'HORS BUDGET',
                }.get(budget_fit, budget_fit)
                table_rows.append([sc_key, f'{sdp_val}m\u00b2', cost_m2_marche, cost_m2_ajuste, cost_total, fit_label])

            # Position: bottom-left, left-aligned under text columns
            # left=0.3", top=3.7", width=5.8", height=1.3"
            tbl_left = Emu(274320)    # 0.3"
            tbl_top = Emu(3383280)    # 3.7"
            tbl_width = Emu(5303520)  # 5.8"
            tbl_height = Emu(1188720) # 1.3"

            rows, cols = 4, 6  # header + 3 data rows
            table_shape = slide17.shapes.add_table(rows, cols, tbl_left, tbl_top, tbl_width, tbl_height)
            tbl = table_shape.table

            # Column widths (proportional)
            col_widths_emu = [
                Emu(914400),   # Scenario (1.0")
                Emu(914400),   # SDP (1.0")
                Emu(1143000),  # Cout/m2 marche (1.25")
                Emu(1143000),  # Cout/m2 ajuste (1.25")
                Emu(1028700),  # Cout total (1.125")
                Emu(1257300),  # Label (1.375")
            ]
            for i, w in enumerate(col_widths_emu):
                tbl.columns[i].width = w

            # Header row
            headers = ['Scenario', 'SDP', 'Cout/m\u00b2\nmarche', 'Cout/m\u00b2\najuste', 'Cout\ntotal', 'Label']
            DARK_GREEN = RGBColor(0x2C, 0x5F, 0x2D)
            WHITE = RGBColor(0xFF, 0xFF, 0xFF)
            LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF0)

            for ci, hdr in enumerate(headers):
                cell = tbl.cell(0, ci)
                cell.text = hdr
                for para in cell.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.CENTER
                    for run in para.runs:
                        run.font.size = Pt(9)
                        run.font.bold = True
                        run.font.color.rgb = WHITE
                # Dark green header background
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_GREEN

            # Data rows
            for ri, row_data in enumerate(table_rows):
                for ci, val in enumerate(row_data):
                    cell = tbl.cell(ri + 1, ci)
                    cell.text = val
                    for para in cell.text_frame.paragraphs:
                        para.alignment = PP_ALIGN.CENTER
                        for run in para.runs:
                            run.font.size = Pt(9)
                            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                    # Alternate row shading
                    if ri % 2 == 1:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = LIGHT_BG

            print("Inserted budget comparison table on slide 17", file=sys.stderr)
        except Exception as e:
            print(f"Warning: could not insert budget table on slide 17: {e}", file=sys.stderr)

    # Slide 18 -- Budget comparison table (same as slide 17, for risk/conclusion)
    if len(slides_list) >= 18:
        slide18 = slides_list[17]
        try:
            table_rows_18 = []
            for sc_key in ['A', 'B', 'C']:
                sdp_val = flat_data.get(f'{sc_key}_sdp', '0')
                cost_m2_marche = flat_data.get(f'{sc_key}_cost_m2_marche', '0k')
                cost_m2_ajuste = flat_data.get(f'{sc_key}_cost_m2_ajuste', '0k')
                cost_total = flat_data.get(f'{sc_key}_cost_total', '0M FCFA')
                budget_fit = flat_data.get(f'{sc_key}_budget_fit', '')
                fit_label = {
                    'DANS_BUDGET': 'DANS BUDGET',
                    'BUDGET_TENDU': 'BUDGET TENDU',
                    'HORS_BUDGET': 'HORS BUDGET',
                }.get(budget_fit, budget_fit)
                table_rows_18.append([sc_key, f'{sdp_val}m\u00b2', cost_m2_marche, cost_m2_ajuste, cost_total, fit_label])

            tbl_left_18 = Emu(1371600)
            tbl_top_18 = Emu(3703320)
            tbl_width_18 = Emu(6400800)
            tbl_height_18 = Emu(1371600)

            table_shape_18 = slide18.shapes.add_table(4, 6, tbl_left_18, tbl_top_18, tbl_width_18, tbl_height_18)
            tbl_18 = table_shape_18.table

            col_widths_18 = [Emu(914400), Emu(914400), Emu(1143000), Emu(1143000), Emu(1028700), Emu(1257300)]
            for i, w in enumerate(col_widths_18):
                tbl_18.columns[i].width = w

            headers_18 = ['Scenario', 'SDP', 'Cout/m\u00b2\nmarche', 'Cout/m\u00b2\najuste', 'Cout\ntotal', 'Label']
            DARK_GREEN_18 = RGBColor(0x2C, 0x5F, 0x2D)
            WHITE_18 = RGBColor(0xFF, 0xFF, 0xFF)
            LIGHT_BG_18 = RGBColor(0xF5, 0xF5, 0xF0)

            for ci, hdr in enumerate(headers_18):
                cell = tbl_18.cell(0, ci)
                cell.text = hdr
                for para in cell.text_frame.paragraphs:
                    para.alignment = PP_ALIGN.CENTER
                    for run in para.runs:
                        run.font.size = Pt(9)
                        run.font.bold = True
                        run.font.color.rgb = WHITE_18
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_GREEN_18

            for ri, row_data in enumerate(table_rows_18):
                for ci, val in enumerate(row_data):
                    cell = tbl_18.cell(ri + 1, ci)
                    cell.text = val
                    for para in cell.text_frame.paragraphs:
                        para.alignment = PP_ALIGN.CENTER
                        for run in para.runs:
                            run.font.size = Pt(9)
                            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                    if ri % 2 == 1:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = LIGHT_BG_18

            print("Inserted budget comparison table on slide 18", file=sys.stderr)
        except Exception as e:
            print(f"Warning: could not insert budget table on slide 18: {e}", file=sys.stderr)

    # Slide 18 -- FALLBACK: Force-inject 3rd column (strategic/phasage) text
    # The template's 3rd column shape has static text instead of a {{invisible_strategic_text}} placeholder,
    # so the placeholder-based replacement never triggers. We find it by position (rightmost text column).
    if len(slides_list) >= 18:
        slide18_fb = slides_list[17]
        strategic_text_s18 = texts.get('invisible_strategic_text_s18', '')
        if strategic_text_s18:
            # Find the rightmost text shape that is NOT the header (header spans full width)
            text_columns = []
            for shape in slide18_fb.shapes:
                if shape.has_text_frame and shape.width < Emu(5000000):  # exclude full-width header
                    text_columns.append(shape)
            # Sort by left position
            text_columns.sort(key=lambda s: s.left)
            if len(text_columns) >= 3:
                third_col = text_columns[2]
                # Check if this shape was NOT already replaced by placeholder logic
                current_text = third_col.text_frame.text
                placeholder_match = re.search(r'\{\{[^}]+\}\}', current_text)
                if not placeholder_match:
                    # Apply the strategic text
                    _clean_text = _clean_phasage_text(strategic_text_s18)
                    _apply_text_to_shape(third_col, '{{invisible_strategic_text}}', _clean_text, 18)
                    print("Slide 18: FALLBACK injected strategic text into 3rd column", file=sys.stderr)
            else:
                print(f"Slide 18: only {len(text_columns)} text columns found (need 3)", file=sys.stderr)
        else:
            print("Slide 18: no invisible_strategic_text_s18 in texts", file=sys.stderr)

    # Slide 19 -- Timeline Gantt chart (below text, full width)
    timeline_chart = chart_paths.get('timeline')
    if timeline_chart and os.path.exists(timeline_chart) and len(slides_list) >= 19:
        slide19 = slides_list[18]
        # left=0.2", top=3.8", width=9.5", height=1.6" → bottom=5.4"
        slide19.shapes.add_picture(timeline_chart,
            Emu(182880), Emu(3474720), Emu(8686800), Emu(1463040))
        print("Inserted timeline chart on slide 19", file=sys.stderr)

    # Slide 20 -- Recap card (below text)
    recap_chart = chart_paths.get('recap_card')
    if recap_chart and os.path.exists(recap_chart) and len(slides_list) >= 20:
        slide20 = slides_list[19]
        # left=0.2", top=4.0", width=9.5", height=1.3" → bottom=5.3"
        slide20.shapes.add_picture(recap_chart,
            Emu(182880), Emu(3657600), Emu(8686800), Emu(1188720))
        print("Inserted recap card on slide 20", file=sys.stderr)

    # -- Final auto-shrink pass for any text-heavy shapes not yet handled --
    # Uses fontScale=80000 (80% minimum) to prevent over-shrinking
    shrink_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_content = shape.text_frame.text.strip()
                if len(text_content) > 150:
                    enable_auto_shrink(shape, fontScale=80000)
                    shrink_count += 1
    print(f"Auto-shrink applied to {shrink_count} text shapes (fontScale=80000)", file=sys.stderr)

    prs.save(output_path)
    print(f"PPTX saved to {output_path}", file=sys.stderr)
    shutil.rmtree(chart_dir, ignore_errors=True)
    shutil.rmtree(img_dir, ignore_errors=True)
    return output_path

if __name__ == '__main__':
    if len(sys.argv) < 4:
        print("Usage: python generate_pptx.py <input.json> <template.pptx> <output.pptx>")
        sys.exit(1)
    input_json = sys.argv[1]
    template_pptx = sys.argv[2]
    output_pptx = sys.argv[3]
    with open(input_json, 'r', encoding='utf-8') as f:
        data = json.load(f)
    assemble_pptx(data, template_pptx, output_pptx)
    print(json.dumps({"status": "ok", "output": output_pptx}))
